# Import library
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from os import remove
import traceback

from common import *

from generate_data import get_statistik_peserta

from pathlib import Path # untuk baca dari file .txt simpan ke variabel

from pptxtopdf import convert # untuk convert PPT ke PDF

def generate_personalized_card(bib_number, nama_peserta, data_path, verbose=True):
    """
    Menghasilkan presentasi PowerPoint (dan mengonversinya ke PDF) yang berisi kartu personalisasi untuk peserta berdasarkan nomor BIB dan nama peserta, menggunakan data dari berbagai grafik dan statistik yang telah dihasilkan sebelumnya.

    Args:
    bib_number (int): Nomor BIB peserta.
    nama_peserta (str): Nama lengkap peserta.
    data_path (dict): Dictionary yang berisi path ke berbagai file grafik dan statistik yang akan dimasukkan ke dalam presentasi.
    """
    try:
        # Variabel-variabel lokal (Path file)
        # Path gambar logo dan sponsor
        # Template
        logo_path = "./img/template/olympiae_logo.png"
        sponsor_path = "./img/template/sponsor_organizer.png"
        sponsor_path_white = "./img/template/sponsor_organizer_white_bg.png"
        image_to_center_path = "./img/template/sample_plot.png"
        motif_batik_path = "./img/template/motifBatik.png"
        motif_batik2_path = "./img/template/motifBatik2.png"
        motif_batik3_path = "./img/template/motifBatik3.png"
        facts_about_you_img_path = "./img/template/facts_about_you.png"
        statistik_peserta_path = "./img/template/your_performance.png"
        statistik_semua_path = "./img/template/community_logo.png"
        closing_image_path = "./img/template/kahf_wardah.png"
        podium_image_path = "./img/template/podium.png"

        # Grafik & Data peserta
        # loading path data
        cumulative_distance_plot_path = data_path["cumulative_distance.png"]
        data_summary_path = data_path["data_summary.png"]
        distance_time_bar_path = data_path["distance_time.png"]
        finish_proportion_path = data_path["finish_proportion.png"]
        finish_per_year_histogram_path = data_path["finish_year_histogram.png"]
        performance_plot_path = data_path["performance_plot.png"]
        scatter_plot_distance_path = data_path["scatter_distance.png"]
        scatter_plot_speed_path = data_path["scatter_speed.png"]
        scatter_plot_elevation_path = data_path["scatter_elevation.png"]
        scatter_plot_elapsed_time_path = data_path["scatter_elapsed_time.png"]
        statistik_individu_path = data_path["statistik_individu.txt"]
        statistik_semua_path = data_path["statistik_semua.txt"]
        statistik_kategori_path = data_path["statistik_kategori.txt"]
        statistik_kegiatan_path = data_path["statistik_jenis_kegiatan.txt"]
        best_stat_path = data_path["statistik_terbaik_individu.txt"]
        facts_about_you_path = data_path["facts_about_you.txt"]

        # Implementasi Fungsi/Prosedur (lokal)
        def baca_file_txt_judul_konten(nama_file):
            """
            membaca file teks, mengambil baris pertama sebagai judul, dan menggabungkan baris-baris sisanya sebagai konten.

            Args:
                nama_file (str): Nama file teks yang akan dibaca.
            Returns:
                Tuple (judul, konten): Judul dan konten dari file teks.
            """
            with open(nama_file, "r", encoding="utf-8") as f:
                lines = f.readlines()  # Membaca semua baris sebagai list

            judul = lines[0].strip() if lines else ""  # Baris pertama sebagai judul
            konten = "".join(lines[1:]).strip() if len(lines) > 1 else ""  # Gabungkan sisa baris sebagai konten
            
            return judul, konten

        def load_konten_teks_PPT(bib_number, nama_peserta):
            """
            Membaca konten teks dari berbagai file .txt dan mengembalikan judul dan konten yang akan digunakan untuk membuat slide PowerPoint.

            Args:
                bib_number (int): Nomor BIB peserta.
                nama_peserta (str): Nama lengkap peserta.
            Returns:
                Tuple: Tuple yang berisi judul dan konten dari berbagai file .txt yang akan digunakan untuk membuat slide PowerPoint.
            """
            try:
                # Fakta tentang kamu
                fact_about_you_content = Path(facts_about_you_path).read_text(encoding="utf-8")
                
                # Best statistics
                best_stat_title, best_stat_content = baca_file_txt_judul_konten(best_stat_path)

                # Statistik Peserta
                statistik_peserta_title, statistik_peserta_content = baca_file_txt_judul_konten(statistik_individu_path)

                # Statistik semua peserta, kategori tertentu, jenis kegiatan tertentu (run/ride)
                statistik_semua_content = Path(statistik_semua_path).read_text(encoding="utf-8")
                statistik_kategori_content = Path(statistik_kategori_path).read_text(encoding="utf-8")
                statistik_kegiatan_content = Path(statistik_kegiatan_path).read_text(encoding="utf-8")

                # Statistik peserta lain
                other_participants_stats_title = "Statistik Peserta Lain"
                other_participants_column_texts = [
                    statistik_semua_content,
                    statistik_kategori_content,
                    statistik_kegiatan_content
                ]

                # Judul slide penutup
                closing_title = "Tukarkan dokumen ini dengan produk Kahf & Wardah saat MudikIAE2025"

                return best_stat_title, best_stat_content, fact_about_you_content, statistik_peserta_title, statistik_peserta_content, other_participants_stats_title, other_participants_column_texts, closing_title
            
            except Exception as e:
                print(e)
                print(traceback.format_exc())  # Print the full traceback (line number & error details)
                if verbose: print("❌Path tidak ditemukan")

        # Fungsi untuk menambahkan logo dan sponsor ke setiap slide
        def add_page_decoration(slide, isbgdark, isSpecial):
            """
            menambahkan dekorasi berupa logo, gambar sponsor, dan motif batik (opsional) ke slide PowerPoint.

            Args:
                slide (Slide): Objek slide PowerPoint yang akan didekorasi.
                isbgdark (bool): Menentukan apakah latar belakang slide gelap.
                isSpecial (bool): Menentukan apakah slide tersebut slide special yang membutuhkan motif batik.
            """
            if isbgdark and isSpecial:
                # slide.shapes.add_picture(motif_batik_path, 0, 0, height=0.5*slide_height)
                slide.shapes.add_picture(motif_batik2_path, 0, Inches(5), width=slide_width)
                slide.shapes.add_picture(motif_batik3_path, 0, 0, width=slide_width)
            slide.shapes.add_picture(logo_path, Inches(8.8), Inches(0), width=Inches(1))
            sponsor_image = sponsor_path if isbgdark else sponsor_path_white
            slide.shapes.add_picture(sponsor_image, 0, Inches(5.65), width=Inches(10))

        # Fungsi untuk membuat slide kosong dengan 1 gambar di tengah : cocok untuk chart
        def add_image_slide(presentation, image_path, isbgdark, image_scale=1, yoffset = 0):
            """
            Menambahkan slide baru dengan gambar yang diposisikan di tengah.

            Args:
                presentation: Objek Presentation dari pptx.
                image_path: Path ke gambar yang akan ditampilkan di tengah slide.
                isbgdark: Boolean untuk menentukan apakah slide memiliki background gelap (True) atau terang (False).
            """
            slide_layout = presentation.slide_layouts[6]  # Blank Slide
            slide = presentation.slides.add_slide(slide_layout)

            # Menentukan ukuran gambar yang proporsional
            chart_width = image_scale*Inches(7)
            chart_height = (1783 / 2441) * chart_width  # Maintain aspect ratio

            # Menambahkan gambar di tengah slide
            slide.shapes.add_picture(image_path,
                                    (slide_width - chart_width) / 2,
                                    (slide_height - chart_height) / 2 - Inches(0.5-yoffset), # yoffset dalam inches
                                    width=chart_width)

            # Menambahkan logo dan sponsor
            add_page_decoration(slide, isbgdark, False)

            return slide  # Jika ingin menyimpan reference ke slide yang dibuat

        # Fungsi untuk menambahkan slide pembuka
        def add_opening_slide(root):
            """
            membuat slide pembuka dengan latar belakang gelap, judul, subjudul (nama peserta), dan dekorasi logo serta sponsor.

            Args:
                root (Presentation): Objek presentasi PowerPoint tempat slide akan ditambahkan.
            """
            first_slide_layout = root.slide_layouts[0]  # Title & Subtitle
            slide1 = root.slides.add_slide(first_slide_layout)

            # Menetapkan warna latar belakang
            background = slide1.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor(4, 25, 77)  # #04194d dalam RGB

            # Menambahkan teks header dengan font elegan dan di tengah
            title = slide1.shapes.title
            title.text = "Thank you for your Participation!"
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Warna putih
            title.text_frame.paragraphs[0].font.size = Pt(54)  # Ukuran font diperbesar
            title.text_frame.paragraphs[0].font.name = "Lucida Calligraphy"  # Font elegan
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Pusatkan teks

            # Pusatkan title secara vertikal
            # title.top = (slide_height - title.height) / 2 - Inches(0.5)

            # Menambahkan subtitle (Nama Peserta)
            subtitle = slide1.placeholders[1]
            subtitle.text = f"{nama_peserta}"
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle.text_frame.paragraphs[0].font.size = Pt(35)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Menambahkan logo dan sponsor
            add_page_decoration(slide1, isbgdark=True, isSpecial=True)

        # Fungsi untuk menambahkan slide berupa satu text box dan sebuah gambar
        def add_text_slide(presentation, text_content, isbgdark=True, image_path=False):
            """
            Menambahkan slide baru dengan teks berformat dalam paragraf.

            Args:
                presentation: Objek Presentation dari pptx.
                text_content: String teks yang akan ditampilkan di slide.
                isbgdark: Boolean untuk menentukan apakah slide memiliki background gelap (True) atau terang (False).
            """
            slide_layout = presentation.slide_layouts[6]  # Blank Slide
            slide = presentation.slides.add_slide(slide_layout)

            # Menetapkan warna latar belakang gelap jika isbgdark=True
            if isbgdark:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(4, 25, 77)  # #04194d dalam RGB

            # Menambahkan teks ke dalam slide
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True  # Agar teks tidak keluar dari kotak
            text_frame.auto_size = True  # Sesuaikan ukuran otomatis

            # Format teks
            for paragraph_text in text_content.split("\n\n"):  # Pisahkan tiap paragraf
                p = text_frame.add_paragraph()
                p.text = paragraph_text
                p.alignment = PP_ALIGN.LEFT  # Teks rata kiri
                p.space_after = Pt(10)  # Jarak antar paragraf

                # Format font
                font = p.font
                font.color.rgb = RGBColor(255, 255, 255)  # Warna putih
                font.size = Pt(10)  # Ukuran teks lebih besar

            # Menambahkan logo dan sponsor
            add_page_decoration(slide, isbgdark, isSpecial=False)

            # menambahkan gambar tambahan
            if (image_path != False):
                slide.shapes.add_picture(image_path, Inches(5), Inches(1), width=Inches(6))

            return slide  # Jika ingin menyimpan reference ke slide yang dibuat

        # Fungsi untuk menambahkan slide berupa judul dan sebuah textbox
        def add_title_paragraph_slide(presentation, title_text, paragraph_text, isbgdark=True, fontSize=20, image_path = False):
            """
            Menambahkan slide dengan title dan paragraf.

            Args:
                presentation: Objek Presentation dari pptx.
                title_text: String untuk judul slide.
                paragraph_text: String untuk isi paragraf.
                isbgdark: Boolean untuk menentukan apakah slide memiliki background gelap (True) atau terang (False).
            """
            slide_layout = presentation.slide_layouts[5]  # Layout Title Only
            slide = presentation.slides.add_slide(slide_layout)

            # Menetapkan warna latar belakang
            if isbgdark:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(4, 25, 77)  # #04194d dalam RGB

            # Menambahkan title
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Format font title
            title_font = title_shape.text_frame.paragraphs[0].font
            title_font.color.rgb = RGBColor(255, 255, 255) if isbgdark else RGBColor(0, 0, 0)
            title_font.size = Pt(36)
            title_font.bold = True

            # Menambahkan paragraf di bawah title
            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for paragraph in paragraph_text.split("\n\n"):  # Pisahkan menjadi paragraf
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.alignment = PP_ALIGN.LEFT  # Rata kiri
                p.space_after = Pt(10)  # Jarak antar paragraf

                # Format font paragraf
                font = p.font
                font.color.rgb = RGBColor(255, 255, 255) if isbgdark else RGBColor(0, 0, 0)
                font.size = Pt(fontSize)

            # Menambahkan logo dan sponsor
            add_page_decoration(slide, isbgdark,isSpecial=False)

            # menambahkan gambar tambahan
            if (image_path != False):
                slide.shapes.add_picture(image_path, Inches(5.5), Inches(2.5), width=Inches(4))


            return slide  # Jika ingin menyimpan reference ke slide yang dibuat

        # Fungsi untuk menambahkan slide berupa tiga kolom teks
        def add_three_column_text_slide(presentation, title_text, column_texts, isbgdark=True):
            """
            Menambahkan slide dengan tiga kolom teks sejajar.

            Args:
                presentation: Objek Presentation dari pptx.
                title_text: String untuk judul slide.
                column_texts: List dari tiga string yang berisi teks untuk masing-masing kolom.
                isbgdark: Boolean untuk menentukan apakah slide memiliki background gelap (True) atau terang (False).
            """
            slide_layout = presentation.slide_layouts[5]  # Layout Title Only
            slide = presentation.slides.add_slide(slide_layout)

            # Menetapkan warna latar belakang
            if isbgdark:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(4, 25, 77)  # #04194d dalam RGB

            # Menambahkan title
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Format font title
            title_font = title_shape.text_frame.paragraphs[0].font
            title_font.color.rgb = RGBColor(255, 255, 255) if isbgdark else RGBColor(0, 0, 0)
            title_font.size = Pt(32)
            title_font.bold = True

            # Posisi untuk tiga kolom
            column_width = Inches(3)
            column_height = Inches(4)
            margin_top = Inches(1)
            col1_x = Inches(0.5)
            col2_x = Inches(3.5)
            col3_x = Inches(6.5)

            # Warna teks berdasarkan latar belakang
            text_color = RGBColor(255, 255, 255) if isbgdark else RGBColor(0, 0, 0)

            # Menambahkan tiga text box sejajar
            for i, (text, x_pos) in enumerate(zip(column_texts, [col1_x, col2_x, col3_x])):
                text_box = slide.shapes.add_textbox(x_pos, margin_top, column_width, column_height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True

                # Tambahkan teks ke dalam textbox
                p = text_frame.add_paragraph()
                p.text = text
                p.alignment = PP_ALIGN.CENTER  # Rata tengah dalam kolom

                # Format font
                font = p.font
                font.color.rgb = text_color
                font.size = Pt(15)

            # Menambahkan logo dan sponsor
            add_page_decoration(slide, isbgdark=True, isSpecial=False)

            return slide  # Jika ingin menyimpan reference ke slide yang dibuat

        # Fungsi untuk menambahkan slide penutup
        def add_closing_slide(presentation, title_text, image_path, isbgdark=True):
            """
            Menambahkan slide penutup dengan teks title dan gambar.

            Args:
                presentation: Objek Presentation dari pptx.
                title_text: String untuk teks judul.
                image_path: Path ke gambar yang akan ditampilkan.
                isbgdark: Boolean untuk menentukan apakah background gelap (True) atau terang (False).
            """
            slide_layout = presentation.slide_layouts[6]  # Blank layout
            slide = presentation.slides.add_slide(slide_layout)

            # Menetapkan warna latar belakang
            if isbgdark:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(4, 25, 77)  # #04194d dalam RGB

            # Menambahkan teks title di tengah atas
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            p = title_frame.add_paragraph()
            p.text = title_text
            p.alignment = PP_ALIGN.CENTER

            # Format font title
            font = p.font
            font.color.rgb = RGBColor(255, 255, 255) if isbgdark else RGBColor(0, 0, 0)
            font.size = Pt(28)
            font.bold = True

            # Menambahkan gambar di tengah slide
            img_width = img_height = Inches(3)
            # img_height = Inches(3.5)  # Sesuaikan aspek rasio jika diperlukan
            slide.shapes.add_picture(image_path, 
                                    (presentation.slide_width - img_width) / 2, 
                                    (presentation.slide_height - img_height) / 2 + Inches(0.5), 
                                    width=img_width)

            # Menambahkan logo dan sponsor
            add_page_decoration(slide, isbgdark=True,isSpecial=True)

            return slide  # Jika ingin menyimpan reference ke slide yang dibuat

        # Membuat objek presentasi
        root = Presentation()
        slide_width = root.slide_width
        slide_height = root.slide_height

        # Loading konten PPT
        best_stat_title, best_stat_content, fact_about_you_content, statistik_peserta_title, statistik_peserta_content, other_participants_stats_title, other_participants_column_texts, closing_title = load_konten_teks_PPT(bib_number, nama_peserta)

        # Slide 1 : Thank you
        add_opening_slide(root)

        # Slide 2 : Statistik peserta
        add_title_paragraph_slide(root, statistik_peserta_title, statistik_peserta_content, isbgdark=True, fontSize=30, image_path=statistik_peserta_path)

        # Slide 3 : Aktivitas Terbaik peserta
        add_title_paragraph_slide(root, best_stat_title, best_stat_content, isbgdark=True, fontSize=12, image_path=podium_image_path)

        # Slide 4: Statistik Keseluruhan Peserta
        add_three_column_text_slide(root, other_participants_stats_title, other_participants_column_texts, isbgdark=True)

        # Slide 5 : Fakta tentang peserta
        add_text_slide(root, fact_about_you_content, isbgdark=True, image_path=facts_about_you_img_path)

        # Slide 6 : Plot Kumulatif Jarak
        add_image_slide(root, cumulative_distance_plot_path , isbgdark=False, image_scale=1.1,yoffset=0.5)

        # Slide 7 : Plot Summary Data
        add_image_slide(root, data_summary_path,isbgdark=False, image_scale=1,yoffset=-0.5)

        # Slide 8 : Plot Jarak terhadap Waktu
        add_image_slide(root, distance_time_bar_path,isbgdark=False, image_scale=1.1,yoffset=0.5)

        # Slide 9 : Proporsi Peserta Finish
        add_image_slide(root, finish_proportion_path,isbgdark=False, image_scale=1)

        # Slide 10 : Histogram Peserta Finish per Tahun
        add_image_slide(root, finish_per_year_histogram_path, isbgdark=False, image_scale=1.2,yoffset=1)

        # Slide 11 : Performance Plot
        add_image_slide(root, performance_plot_path,isbgdark=False, image_scale=1.1)

        # Slide 12 : Scatter Plot Jarak
        add_image_slide(root, scatter_plot_distance_path, isbgdark=False, image_scale=1,yoffset=-0.5)

        # Slide 13 : Scatter Plot Elapsed Time
        add_image_slide(root, scatter_plot_elapsed_time_path, isbgdark=False, image_scale=1,yoffset=-0.5)

        # Slide 14 : Scatter Plot Elevation
        add_image_slide(root, scatter_plot_elevation_path, isbgdark=False, image_scale=1,yoffset=-0.5)

        # Slide 15 : Scatter Plot Speed
        add_image_slide(root, scatter_plot_speed_path, isbgdark=False, image_scale=1,yoffset=-0.5)

        # Slide 16 : Undangan ke MudikIAE (Slide penutup)
        add_closing_slide(root, closing_title, closing_image_path, isbgdark=True)

        # Menyimpan file
        root.save(f"./output/{bib_number:04d}_{nama_peserta}/{bib_number:04d}_Personalized_Card.pptx")

        # PPT to PDF
        input_dir = f"./output/{bib_number:04d}_{nama_peserta}/{bib_number:04d}_Personalized_Card.pptx"
        output_dir = f"./output/{bib_number:04d}_{nama_peserta}/"

        convert(input_dir, output_dir)
        remove(input_dir) # menghapus file PPT

        if verbose: print("✅ Proses generate personalized card berhasil")

    except Exception as e:
        print(e)
        print(traceback.format_exc())  # Print the full traceback (line number & error details)
        if verbose: print("❌ Proses generate personalized card gagal")
